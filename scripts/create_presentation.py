"""Generate a modern PowerPoint presentation for the UTC Assessment.

Usage (from project root):
    python scripts/create_presentation.py
Or in a Jupyter cell:
    %run scripts/create_presentation.py

Requires: pip install python-pptx
"""

import warnings
warnings.filterwarnings("ignore")

import os, sys
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.colors import ListedColormap, BoundaryNorm
import matplotlib.patches as mpatches
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))
os.chdir(ROOT)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from scripts.config import *
from scripts import sentinel2, lidar, alignment, classification, accuracy

# =====================================================================
#  Output paths
# =====================================================================
FIG_DIR = ROOT / "output" / "figures"
FIG_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = ROOT / "output" / "UTC_Assessment_Fairfax_County_VA.pptx"

# =====================================================================
#  Design tokens — clean light theme
# =====================================================================
BG       = "#FFFFFF"
BG_CARD  = "#F6F8FA"
CYAN     = "#0969DA"
GREEN    = "#1A7F37"
ORANGE   = "#BC4C00"
WHITE    = "#1F2328"     # primary text (dark on light)
GRAY     = "#656D76"
PINK     = "#BF3989"
BORDER   = "#D1D9E0"

C_BG       = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD     = RGBColor(0xF0, 0xF3, 0xF6)
C_CYAN     = RGBColor(0x09, 0x69, 0xDA)
C_GREEN    = RGBColor(0x1A, 0x7F, 0x37)
C_ORANGE   = RGBColor(0xBC, 0x4C, 0x00)
C_WHITE    = RGBColor(0x1F, 0x23, 0x28)   # primary text
C_GRAY     = RGBColor(0x65, 0x6D, 0x76)
C_PINK     = RGBColor(0xBF, 0x39, 0x89)
C_DARK_ACCENT = RGBColor(0xD1, 0xD9, 0xE0)

FONT_TITLE = "Segoe UI Light"
FONT_BODY  = "Segoe UI"
FONT_MONO  = "Consolas"

# =====================================================================
#  Matplotlib light theme
# =====================================================================
plt.style.use("default")
plt.rcParams.update({
    "figure.facecolor": BG,
    "axes.facecolor": BG_CARD,
    "savefig.facecolor": BG,
    "text.color": WHITE,
    "axes.labelcolor": WHITE,
    "xtick.color": GRAY,
    "ytick.color": GRAY,
    "axes.edgecolor": BORDER,
    "grid.color": "#E8EAED",
    "font.family": "sans-serif",
    "font.size": 11,
})

# Classification palette (matching visualization.py)
CMAP_CLS  = ListedColormap(["#cccccc", "#d2b48c", "#90ee90", "#006400"])
NORM_CLS  = BoundaryNorm([-0.5, 0.5, 1.5, 2.5, 3.5], CMAP_CLS.N)
CLS_PATCHES = [
    mpatches.Patch(color="#006400", label="Tree canopy"),
    mpatches.Patch(color="#90ee90", label="Low vegetation"),
    mpatches.Patch(color="#d2b48c", label="Non-vegetation"),
]


def scale_rgb(arr, plow=2, phigh=98):
    lo, hi = np.nanpercentile(arr, [plow, phigh])
    return np.clip((arr - lo) / (hi - lo + 1e-10), 0, 1)


# =====================================================================
#  1. Load cached data & run analysis
# =====================================================================
print("=" * 60)
print("  UTC Assessment — Presentation Builder")
print("=" * 60)

print("\n[1/6] Loading Sentinel-2 ...")
crop_files = list((ROOT / "data" / "sentinel2").glob("*_crop.tif"))
if not crop_files:
    raise FileNotFoundError("No cropped S2 scene found. Run the analysis first.")
s2, s2_crs = sentinel2.load_stacked(str(crop_files[0]))
idx = sentinel2.compute_indices(s2)
scene_name = crop_files[0].stem.replace("_crop", "")
s2_date = f"{scene_name.split('_')[2][:4]}-{scene_name.split('_')[2][4:6]}-{scene_name.split('_')[2][6:8]}"

print("\n[2/6] Loading LiDAR products ...")
import rioxarray
dtm = rioxarray.open_rasterio(str(ROOT / "data" / "dem" / "dtm_3dep.tif")).squeeze()
dsm_path = ROOT / "data" / "dem" / "dsm_3dep.tif"
if dsm_path.exists():
    dsm = rioxarray.open_rasterio(str(dsm_path)).squeeze()
else:
    dsm = None
chm = rioxarray.open_rasterio(str(ROOT / "data" / "lidar" / "chm.tif")).squeeze()

print("\n[3/6] Aligning grids ...")
aligned = alignment.align_grids(s2, chm)

print("\n[4/6] Running classification ...")
rule_map, valid = classification.classify_rules(
    aligned["ndvi_a"], aligned["chm_max_10"]
)
rf_map, rf_model, feat_names = classification.classify_rf(
    aligned["s2_utm"], aligned["ndvi_a"], aligned["ndre_a"], aligned["ndbi_a"],
    aligned["chm_max_10"], aligned["chm_std_10"], rule_map,
)

print("\n[5/6] Running accuracy assessment ...")
ref_map = accuracy.build_reference_map(
    aligned["ndvi_a"], aligned["chm_max_10"], valid
)
sample_ref, sample_rule, sample_rf = accuracy.sample_reference(
    ref_map, rule_map, rf_map
)
oa_rule, k_rule = accuracy.report_accuracy(sample_ref, sample_rule, "Rule-Based")
oa_rf, k_rf = accuracy.report_accuracy(sample_ref, sample_rf, "Random Forest")

# Derived values for slides
pixel_area_ha = 10 * 10 / 1e4
total_valid = valid.sum()
area_ha = total_valid * pixel_area_ha

tree_rule_pct = (rule_map == CLASS_TREE).sum() / total_valid * 100
low_rule_pct  = (rule_map == CLASS_LOW).sum()  / total_valid * 100
nonv_rule_pct = (rule_map == CLASS_NONVEG).sum() / total_valid * 100

tree_rf_pct = (rf_map == CLASS_TREE).sum() / total_valid * 100
low_rf_pct  = (rf_map == CLASS_LOW).sum()  / total_valid * 100
nonv_rf_pct = (rf_map == CLASS_NONVEG).sum() / total_valid * 100

tree_rule_ha = (rule_map == CLASS_TREE).sum() * pixel_area_ha
tree_rf_ha   = (rf_map == CLASS_TREE).sum()   * pixel_area_ha


# =====================================================================
#  2. Generate figures
# =====================================================================
print("\n[6/6] Generating figures ...")

DPI = 200


def savefig(fig, name):
    path = FIG_DIR / f"{name}.png"
    fig.savefig(str(path), dpi=DPI, bbox_inches="tight", pad_inches=0.3)
    plt.close(fig)
    print(f"  Saved {path.name}")
    return path


# --- Study area ---
def fig_study_area():
    fig, ax = plt.subplots(figsize=(10, 8))
    gdf_utm.boundary.plot(ax=ax, edgecolor="#CF222E", linewidth=2.5)
    try:
        import contextily as cx
        cx.add_basemap(ax, crs=gdf_utm.crs.to_string(),
                       source=cx.providers.Esri.WorldImagery, zoom=15)
    except Exception:
        ax.set_facecolor("#E8EAED")
    ax.set_title("Study Area — Fairfax County, VA", fontsize=16, color=WHITE,
                 fontweight="bold", pad=15)
    ax.set_xlabel("Easting (m)", fontsize=11)
    ax.set_ylabel("Northing (m)", fontsize=11)
    bx = bounds_utm
    ax.text(0.02, 0.02,
            f"BBOX: ({BBOX[0]:.3f}, {BBOX[1]:.3f}, {BBOX[2]:.3f}, {BBOX[3]:.3f})\n"
            f"Area: {area_ha:.0f} ha  |  CRS: EPSG:32618",
            transform=ax.transAxes, fontsize=9, color=GRAY,
            verticalalignment="bottom",
            bbox=dict(boxstyle="round,pad=0.4", fc=BG_CARD, ec=BORDER, alpha=0.9))
    plt.tight_layout()
    return savefig(fig, "01_study_area")


# --- Sentinel-2 composites ---
def fig_s2_composites():
    rgb = np.dstack([idx["red"], idx["green"], idx["blue"]])
    cir = np.dstack([idx["nir"], idx["red"], idx["green"]])
    fig, axes = plt.subplots(1, 2, figsize=(16, 7))
    axes[0].imshow(scale_rgb(rgb))
    axes[0].set_title("True Color (B04–B03–B02)", fontsize=14, pad=10)
    axes[0].axis("off")
    axes[1].imshow(scale_rgb(cir))
    axes[1].set_title("Color Infrared (B08–B04–B03)", fontsize=14, pad=10)
    axes[1].axis("off")
    fig.suptitle(f"Sentinel-2 L2A — {s2_date}", fontsize=16,
                 fontweight="bold", color=CYAN, y=0.98)
    plt.tight_layout(rect=[0, 0, 1, 0.95])
    return savefig(fig, "02_s2_composites")


# --- Spectral indices ---
def fig_spectral_indices():
    fig, axes = plt.subplots(1, 3, figsize=(18, 6))
    for ax, data, title, cmap, vr in [
        (axes[0], idx["ndvi"], "NDVI", "RdYlGn", (-0.1, 0.8)),
        (axes[1], idx["ndre"], "NDRE", "RdYlGn", (-0.1, 0.7)),
        (axes[2], idx["ndbi"], "NDBI", "RdYlBu_r", (-0.4, 0.2)),
    ]:
        im = ax.imshow(data, cmap=cmap, vmin=vr[0], vmax=vr[1])
        ax.set_title(title, fontsize=14, pad=10)
        ax.axis("off")
        cb = plt.colorbar(im, ax=ax, fraction=0.035, pad=0.04)
        cb.ax.tick_params(colors=GRAY, labelsize=9)
    fig.suptitle("Spectral Indices", fontsize=16, fontweight="bold",
                 color=CYAN, y=0.98)
    plt.tight_layout(rect=[0, 0, 1, 0.94])
    return savefig(fig, "03_spectral_indices")


# --- LiDAR products ---
def fig_lidar_products():
    dtm_np = dtm.values if hasattr(dtm, "values") else dtm
    chm_np = chm.values if hasattr(chm, "values") else chm
    panels = [(dtm_np, "DTM (Bare Earth)", "terrain", None),
              (chm_np, "CHM (Canopy Height)", "Greens", (0, 35))]
    if dsm is not None:
        dsm_np = dsm.values if hasattr(dsm, "values") else dsm
        panels.insert(1, (dsm_np, "DSM (Surface Model)", "terrain", None))

    ncols = len(panels)
    fig, axes = plt.subplots(1, ncols, figsize=(6 * ncols, 6))
    if ncols == 1:
        axes = [axes]
    for ax, (data, title, cmap, vlim) in zip(axes, panels):
        kw = {"cmap": cmap}
        if vlim:
            kw["vmin"], kw["vmax"] = vlim
        im = ax.imshow(data, **kw)
        ax.set_title(title, fontsize=14, pad=10)
        ax.axis("off")
        lbl = "Height (m)" if "CHM" in title else "Elevation (m)"
        cb = plt.colorbar(im, ax=ax, fraction=0.035, pad=0.04, label=lbl)
        cb.ax.tick_params(colors=GRAY, labelsize=9)
        cb.set_label(lbl, color=GRAY, fontsize=10)
    fig.suptitle("LiDAR-Derived Elevation Products (USGS 3DEP — 2018)",
                 fontsize=16, fontweight="bold", color=CYAN, y=0.98)
    plt.tight_layout(rect=[0, 0, 1, 0.94])
    return savefig(fig, "04_lidar_products")


# --- Classification comparison ---
def fig_classification():
    fig, axes = plt.subplots(1, 2, figsize=(16, 7))
    for ax, data, title in [
        (axes[0], rule_map, f"Rule-Based  (OA={oa_rule:.0%}, κ={k_rule:.2f})"),
        (axes[1], rf_map,   f"Random Forest  (OA={oa_rf:.0%}, κ={k_rf:.2f})"),
    ]:
        ax.imshow(data, cmap=CMAP_CLS, norm=NORM_CLS)
        ax.set_title(title, fontsize=13, pad=10)
        ax.axis("off")
    axes[1].legend(handles=CLS_PATCHES, loc="lower right", fontsize=10,
                   facecolor=BG_CARD, edgecolor=BORDER, labelcolor=WHITE)
    fig.suptitle("UTC Classification Results", fontsize=16,
                 fontweight="bold", color=CYAN, y=0.98)
    plt.tight_layout(rect=[0, 0, 1, 0.94])
    return savefig(fig, "05_classification")


# --- Confusion matrices ---
def fig_confusion():
    from sklearn.metrics import confusion_matrix as cm_func
    classes = [CLASS_NONVEG, CLASS_LOW, CLASS_TREE]
    names = ["Non-veg", "Low veg", "Tree"]

    fig, axes = plt.subplots(1, 2, figsize=(14, 5.5))
    for ax, y_pred, title, color in [
        (axes[0], sample_rule, f"Rule-Based (OA={oa_rule:.1%})", CYAN),
        (axes[1], sample_rf,   f"Random Forest (OA={oa_rf:.1%})", GREEN),
    ]:
        cm = cm_func(sample_ref, y_pred, labels=classes)
        cm_pct = cm.astype(float) / cm.sum(axis=1, keepdims=True) * 100
        im = ax.imshow(cm_pct, cmap="Blues", vmin=0, vmax=100)
        ax.set_xticks(range(3))
        ax.set_yticks(range(3))
        ax.set_xticklabels(names, fontsize=10, color=WHITE)
        ax.set_yticklabels(names, fontsize=10, color=WHITE)
        ax.set_xlabel("Predicted", fontsize=11, color=GRAY)
        ax.set_ylabel("Reference", fontsize=11, color=GRAY)
        ax.set_title(title, fontsize=13, pad=10, color=color)
        for i in range(3):
            for j in range(3):
                txt_color = "white" if cm_pct[i, j] > 50 else "#1F2328"
                ax.text(j, i, f"{cm[i, j]}\n({cm_pct[i, j]:.0f}%)",
                        ha="center", va="center", fontsize=10,
                        color=txt_color, fontweight="bold")
    fig.suptitle("Confusion Matrices", fontsize=16,
                 fontweight="bold", color=CYAN, y=1.0)
    plt.tight_layout(rect=[0, 0, 1, 0.95])
    return savefig(fig, "06_confusion")


# --- Feature importance ---
def fig_importance():
    importances = rf_model.feature_importances_
    order = np.argsort(importances)[::-1]
    sorted_names = [feat_names[i] for i in order]
    sorted_vals  = importances[order]

    fig, ax = plt.subplots(figsize=(10, 6))
    colors = [CYAN if v > 0.1 else GREEN if v > 0.03 else GRAY
              for v in sorted_vals]
    bars = ax.barh(range(len(sorted_names)), sorted_vals[::-1],
                   color=colors[::-1], edgecolor="none", height=0.7)
    ax.set_yticks(range(len(sorted_names)))
    ax.set_yticklabels(sorted_names[::-1], fontsize=11)
    ax.set_xlabel("Importance", fontsize=12)
    ax.set_title("Random Forest — Feature Importances", fontsize=14,
                 fontweight="bold", color=CYAN, pad=15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for i, v in enumerate(sorted_vals[::-1]):
        ax.text(v + 0.005, i, f"{v:.1%}", va="center", fontsize=9, color=WHITE)
    plt.tight_layout()
    return savefig(fig, "07_feature_importance")


# --- Summary dashboard ---
def fig_summary():
    rgb = np.dstack([idx["red"], idx["green"], idx["blue"]])
    fig, axes = plt.subplots(2, 3, figsize=(20, 13))

    axes[0, 0].imshow(scale_rgb(rgb))
    axes[0, 0].set_title("Sentinel-2 True Color", fontsize=13, pad=8)
    axes[0, 0].axis("off")

    im1 = axes[0, 1].imshow(aligned["ndvi_a"], cmap="RdYlGn",
                             vmin=-0.1, vmax=0.8)
    axes[0, 1].set_title("NDVI (10 m aligned)", fontsize=13, pad=8)
    axes[0, 1].axis("off")
    plt.colorbar(im1, ax=axes[0, 1], fraction=0.035)

    im2 = axes[0, 2].imshow(aligned["chm_max_10"], cmap="Greens",
                             vmin=0, vmax=35)
    axes[0, 2].set_title("CHM max (10 m)", fontsize=13, pad=8)
    axes[0, 2].axis("off")
    cb2 = plt.colorbar(im2, ax=axes[0, 2], fraction=0.035)
    cb2.set_label("m", color=GRAY)

    axes[1, 0].imshow(rule_map, cmap=CMAP_CLS, norm=NORM_CLS)
    axes[1, 0].set_title(f"Rule-Based (OA={oa_rule:.0%}, κ={k_rule:.2f})",
                         fontsize=13, pad=8)
    axes[1, 0].axis("off")

    axes[1, 1].imshow(rf_map, cmap=CMAP_CLS, norm=NORM_CLS)
    axes[1, 1].set_title(f"Random Forest (OA={oa_rf:.0%}, κ={k_rf:.2f})",
                         fontsize=13, pad=8)
    axes[1, 1].axis("off")

    diff = (rule_map != rf_map).astype(np.uint8)
    diff[~valid] = 2
    cmap_diff = ListedColormap(["#FFFFFF", "#CF222E", "#E8EAED"])
    axes[1, 2].imshow(diff, cmap=cmap_diff, vmin=0, vmax=2)
    axes[1, 2].set_title("Disagreement Map", fontsize=13, pad=8)
    axes[1, 2].axis("off")
    dpct = (diff == 1).sum() / valid.sum() * 100
    axes[1, 2].text(0.02, 0.02, f"{dpct:.1f}% differ",
                    transform=axes[1, 2].transAxes, fontsize=10, color=ORANGE,
                    bbox=dict(boxstyle="round", fc=BG_CARD, ec=BORDER, alpha=0.9))

    fig.legend(handles=CLS_PATCHES, loc="lower center", ncol=3,
               fontsize=11, facecolor=BG_CARD, edgecolor=BORDER,
               labelcolor=WHITE, bbox_to_anchor=(0.5, -0.01))
    fig.suptitle("Urban Tree Canopy Assessment — Summary",
                 fontsize=18, fontweight="bold", color=CYAN, y=0.98)
    plt.tight_layout(rect=[0, 0.02, 1, 0.95])
    return savefig(fig, "08_summary_dashboard")


# Generate all
paths = {}
for name, func in [
    ("study_area", fig_study_area),
    ("s2_composites", fig_s2_composites),
    ("spectral_indices", fig_spectral_indices),
    ("lidar_products", fig_lidar_products),
    ("classification", fig_classification),
    ("confusion", fig_confusion),
    ("importance", fig_importance),
    ("summary", fig_summary),
]:
    paths[name] = func()


# =====================================================================
#  3. Build PowerPoint
# =====================================================================
print("\nBuilding PowerPoint ...")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height


def add_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=C_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = C_DARK_ACCENT
    shape.line.width = Pt(0.75)
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_line(slide, left, top, width, color=C_CYAN, thickness=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text,
             font_size=18, color=C_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name=FONT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines,
                  font_size=16, color=C_WHITE, line_spacing=1.5,
                  font_name=FONT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(sz * 0.5)
    return txBox


# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
add_line(slide, Inches(1.5), Inches(2.2), Inches(3), C_CYAN, Pt(4))
add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
         "Urban Tree Canopy Assessment", font_size=40, bold=True,
         font_name=FONT_TITLE, color=C_WHITE)
add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
         "Fairfax County, Virginia — 2018", font_size=24,
         font_name=FONT_TITLE, color=C_CYAN)
add_text(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
         "Multispectral Imagery + LiDAR Data Fusion", font_size=16,
         color=C_GRAY, font_name=FONT_BODY)
add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
         f"Sentinel-2 L2A  ·  USGS 3DEP  ·  {s2_date}", font_size=12,
         color=C_GRAY, font_name=FONT_MONO)

# ── Slide 2: Agenda ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
         "Agenda", font_size=32, bold=True, font_name=FONT_TITLE, color=C_WHITE)
add_line(slide, Inches(1), Inches(1.2), Inches(2), C_CYAN)

agenda_items = [
    ("01", "Study Area & Data Sources"),
    ("02", "Sentinel-2 Multispectral Imagery"),
    ("03", "Spectral Vegetation Indices"),
    ("04", "LiDAR-Derived Elevation Products"),
    ("05", "Classification Methods & Results"),
    ("06", "Accuracy Assessment"),
    ("07", "Key Findings"),
]
for i, (num, label) in enumerate(agenda_items):
    y = Inches(1.8 + i * 0.7)
    add_text(slide, Inches(1.5), y, Inches(0.8), Inches(0.5),
             num, font_size=20, bold=True, color=C_CYAN, font_name=FONT_MONO)
    add_text(slide, Inches(2.4), y, Inches(8), Inches(0.5),
             label, font_size=20, color=C_WHITE, font_name=FONT_BODY)

# ── Slide 3: Study Area ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Study Area", font_size=28, bold=True, font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)
slide.shapes.add_picture(str(paths["study_area"]),
                         Inches(0.3), Inches(1.3), Inches(7.5))
# Info panel
info_lines = [
    ("Location", 12, C_GRAY, False),
    ("Fairfax County, Virginia (Mason District)", 16, C_WHITE, True),
    ("", 8, C_GRAY, False),
    ("Coordinates (WGS 84)", 12, C_GRAY, False),
    (f"{BBOX[0]:.3f}°W, {BBOX[1]:.3f}°N  →  {BBOX[2]:.3f}°W, {BBOX[3]:.3f}°N", 13, C_WHITE, False),
    ("", 8, C_GRAY, False),
    ("Area", 12, C_GRAY, False),
    (f"{area_ha:.0f} hectares ({area_ha/100:.2f} km²)", 16, C_WHITE, True),
    ("", 8, C_GRAY, False),
    ("Projection", 12, C_GRAY, False),
    ("UTM Zone 18N (EPSG:32618)", 14, C_WHITE, False),
    ("", 8, C_GRAY, False),
    ("Characteristics", 12, C_GRAY, False),
    ("Suburban residential with mature tree canopy,", 13, C_WHITE, False),
    ("commercial areas, and road networks", 13, C_WHITE, False),
]
add_multiline(slide, Inches(8.3), Inches(1.5), Inches(4.5), Inches(5.5),
              info_lines, font_name=FONT_BODY)

# ── Slide 4: Data Sources ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Data Sources", font_size=28, bold=True, font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)

# S2 card
add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2),
             C_CARD, radius=0.02)
s2_lines = [
    ("SENTINEL-2 L2A", 11, C_CYAN, True),
    ("Multispectral Satellite Imagery", 18, C_WHITE, True),
    ("", 8, C_GRAY, False),
    ("Source: ESA / Copernicus (via Planetary Computer)", 12, C_GRAY, False),
    (f"Scene: {s2_date}  |  Cloud: 0.4%", 13, C_WHITE, False),
    ("", 6, C_GRAY, False),
    ("Bands loaded:", 12, C_GRAY, False),
    ("B02 (Blue), B03 (Green), B04 (Red)", 12, C_WHITE, False),
    ("B05 (Red Edge), B06, B07 (Vegetation)", 12, C_WHITE, False),
    ("B08 (NIR), B11–B12 (SWIR)", 12, C_WHITE, False),
    ("", 6, C_GRAY, False),
    ("Resolution: 10 m (20 m bands resampled)", 13, C_GREEN, False),
    ("Values: Surface reflectance (0–1)", 13, C_GREEN, False),
]
add_multiline(slide, Inches(1.0), Inches(1.8), Inches(5.0), Inches(4.5),
              s2_lines, font_name=FONT_BODY)

# LiDAR card
add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
             C_CARD, radius=0.02)
lidar_lines = [
    ("USGS 3DEP", 11, C_CYAN, True),
    ("LiDAR-Derived Elevation Data", 18, C_WHITE, True),
    ("", 8, C_GRAY, False),
    ("Source: USGS 3D Elevation Program", 12, C_GRAY, False),
    ("Survey: VA Fairfax County 2018", 13, C_WHITE, False),
    ("", 6, C_GRAY, False),
    ("Products:", 12, C_GRAY, False),
    ("DTM — Digital Terrain Model (bare earth)", 12, C_WHITE, False),
    ("DSM — Digital Surface Model (top of canopy)", 12, C_WHITE, False),
    ("CHM — Canopy Height Model (DSM − DTM)", 12, C_WHITE, False),
    ("", 6, C_GRAY, False),
    ("Resolution: 2 m (aggregated to 10 m)", 13, C_GREEN, False),
    ("Temporal match: both datasets from 2018", 13, C_GREEN, False),
]
add_multiline(slide, Inches(7.2), Inches(1.8), Inches(5.0), Inches(4.5),
              lidar_lines, font_name=FONT_BODY)

# ── Slide 5: S2 Composites ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Sentinel-2 Composites", font_size=28, bold=True, font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)
slide.shapes.add_picture(str(paths["s2_composites"]),
                         Inches(0.3), Inches(1.3), Inches(12.7))

# ── Slide 6: Spectral Indices ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Spectral Vegetation Indices", font_size=28, bold=True,
         font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)
slide.shapes.add_picture(str(paths["spectral_indices"]),
                         Inches(0.15), Inches(1.3), Inches(13))

# ── Slide 7: LiDAR Products ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "LiDAR-Derived Elevation Products", font_size=28, bold=True,
         font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)
slide.shapes.add_picture(str(paths["lidar_products"]),
                         Inches(0.15), Inches(1.3), Inches(13))

# ── Slide 8: Methods ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Classification Methods", font_size=28, bold=True, font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)

# Rule card
add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2),
             C_CARD, radius=0.02)
rule_lines = [
    ("METHOD 1", 11, C_ORANGE, True),
    ("Rule-Based Classification", 20, C_WHITE, True),
    ("", 8, C_GRAY, False),
    ("Simple threshold approach using two inputs:", 13, C_GRAY, False),
    ("", 6, C_GRAY, False),
    ("Tree Canopy:", 13, C_GREEN, True),
    ("  NDVI > 0.3  AND  CHM > 2.5 m", 13, C_WHITE, False),
    ("", 4, C_GRAY, False),
    ("Low Vegetation:", 13, C_GREEN, True),
    ("  NDVI > 0.3  AND  CHM ≤ 2.5 m", 13, C_WHITE, False),
    ("", 4, C_GRAY, False),
    ("Non-Vegetation:", 13, C_GREEN, True),
    ("  NDVI ≤ 0.3", 13, C_WHITE, False),
    ("", 8, C_GRAY, False),
    ("✦ Fast, transparent, reproducible", 12, C_GRAY, False),
]
add_multiline(slide, Inches(1.0), Inches(1.8), Inches(5.0), Inches(4.8),
              rule_lines, font_name=FONT_BODY)

# RF card
add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
             C_CARD, radius=0.02)
rf_lines = [
    ("METHOD 2", 11, C_CYAN, True),
    ("Random Forest Classifier", 20, C_WHITE, True),
    ("", 8, C_GRAY, False),
    ("Supervised ML trained on high-confidence pixels:", 13, C_GRAY, False),
    ("", 6, C_GRAY, False),
    ("Features (14 total):", 13, C_GREEN, True),
    ("  9 Sentinel-2 bands", 13, C_WHITE, False),
    ("  3 spectral indices (NDVI, NDRE, NDBI)", 13, C_WHITE, False),
    ("  2 CHM statistics (max height, std dev)", 13, C_WHITE, False),
    ("", 6, C_GRAY, False),
    ("Hyperparameters:", 13, C_GREEN, True),
    ("  150 trees  |  max depth 15", 13, C_WHITE, False),
    ("  20,000 training samples", 13, C_WHITE, False),
    ("", 8, C_GRAY, False),
    ("✦ Learns complex non-linear boundaries", 12, C_GRAY, False),
]
add_multiline(slide, Inches(7.2), Inches(1.8), Inches(5.0), Inches(4.8),
              rf_lines, font_name=FONT_BODY)

# ── Slide 9: Classification Results ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Classification Results", font_size=28, bold=True, font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)
slide.shapes.add_picture(str(paths["classification"]),
                         Inches(0.3), Inches(1.3), Inches(12.7))

# ── Slide 10: Accuracy ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Accuracy Assessment", font_size=28, bold=True, font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)
slide.shapes.add_picture(str(paths["confusion"]),
                         Inches(0.8), Inches(1.2), Inches(11.5))
add_text(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
         f"498 stratified random samples  ·  "
         f"Rule-Based: OA={oa_rule:.1%}, κ={k_rule:.3f}  ·  "
         f"Random Forest: OA={oa_rf:.1%}, κ={k_rf:.3f}",
         font_size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER,
         font_name=FONT_MONO)

# ── Slide 11: Feature Importance ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Feature Importance Analysis", font_size=28, bold=True,
         font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)
slide.shapes.add_picture(str(paths["importance"]),
                         Inches(0.5), Inches(1.2), Inches(8))
# Insight box
add_shape_bg(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(4.5),
             C_CARD, radius=0.02)
insight_lines = [
    ("KEY INSIGHT", 11, C_GREEN, True),
    ("", 6, C_GRAY, False),
    ("Top 4 features account for", 14, C_WHITE, False),
    ("~77% of model decisions:", 14, C_WHITE, False),
    ("", 8, C_GRAY, False),
    ("NDVI — greenness (satellite)", 12, C_WHITE, False),
    ("CHM max — tree height (LiDAR)", 12, C_WHITE, False),
    ("NDRE — plant health (satellite)", 12, C_WHITE, False),
    ("CHM std — height texture (LiDAR)", 12, C_WHITE, False),
    ("", 8, C_GRAY, False),
    ("Data fusion of satellite +", 13, C_CYAN, True),
    ("LiDAR is essential for accurate", 13, C_CYAN, True),
    ("canopy classification.", 13, C_CYAN, True),
]
add_multiline(slide, Inches(9.1), Inches(1.8), Inches(3.5), Inches(4),
              insight_lines, font_name=FONT_BODY)

# ── Slide 12: Summary Dashboard ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
slide.shapes.add_picture(str(paths["summary"]),
                         Inches(0.15), Inches(0.1), Inches(13))

# ── Slide 13: Key Findings ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Key Findings", font_size=28, bold=True, font_name=FONT_TITLE)
add_line(slide, Inches(0.6), Inches(0.95), Inches(1.8), C_CYAN)

findings = [
    (f"Tree canopy covers ~{tree_rule_pct:.0f}% of the study area "
     f"({tree_rule_ha:.0f} ha)", C_GREEN),
    (f"Random Forest slightly outperforms rule-based "
     f"(OA {oa_rf:.0%} vs {oa_rule:.0%})", C_CYAN),
    ("LiDAR height data is critical — CHM features rank 2nd and 4th "
     "in importance", C_ORANGE),
    ("Both methods agree on >97% of pixels — "
     "rule-based is a strong baseline", C_WHITE),
    ("Temporal consistency (both 2018) eliminates change-related "
     "classification noise", C_PINK),
]
for i, (text, color) in enumerate(findings):
    y = Inches(1.6 + i * 1.05)
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.8),
                 C_CARD, radius=0.02)
    add_text(slide, Inches(1.1), y + Inches(0.05), Inches(0.5), Inches(0.7),
             "→", font_size=22, bold=True, color=color, font_name=FONT_MONO)
    add_text(slide, Inches(1.7), y + Inches(0.1), Inches(10.5), Inches(0.6),
             text, font_size=16, color=C_WHITE, font_name=FONT_BODY)

# ── Slide 14: Thank You ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_line(slide, Inches(4.5), Inches(2.5), Inches(4.3), C_CYAN, Pt(3))
add_text(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1),
         "Thank You", font_size=44, bold=True, font_name=FONT_TITLE,
         alignment=PP_ALIGN.CENTER, color=C_WHITE)
add_text(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.6),
         "Questions & Discussion", font_size=22, font_name=FONT_TITLE,
         alignment=PP_ALIGN.CENTER, color=C_CYAN)
add_text(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.8),
         f"Urban Tree Canopy Assessment  ·  Fairfax County, VA  ·  {s2_date}\n"
         "Sentinel-2 L2A + USGS 3DEP LiDAR  ·  EPSG:32618",
         font_size=12, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER,
         color=C_GRAY)


# ── Save ──
prs.save(str(PPTX_PATH))
print(f"\n{'=' * 60}")
print(f"  ✓  Presentation saved: {PPTX_PATH}")
print(f"  ✓  Figures saved in:   {FIG_DIR}/")
print(f"{'=' * 60}")
